import os
import logging
from typing import List, Dict, Callable, Optional, Any, Union
from pathlib import Path

# Configure logging
logger = logging.getLogger(__name__)

# --- Modular File Loader & Data Connector Registry ---

FILE_LOADERS = {}
DATA_CONNECTORS = {}

def register_file_loader(extension: str, loader_fn: Callable[[str], Optional[str]]):
    """Register a file loader for a given extension."""
    FILE_LOADERS[extension.lower()] = loader_fn

def get_file_loader(extension: str) -> Optional[Callable[[str], Optional[str]]]:
    return FILE_LOADERS.get(extension.lower())

def register_data_connector(source_type: str, connector_fn: Callable[..., List[Dict[str, Any]]]):
    """Register a data connector for a given source type (e.g., mysql, notion, sharepoint)."""
    DATA_CONNECTORS[source_type.lower()] = connector_fn

def get_data_connector(source_type: str) -> Optional[Callable[..., List[Dict[str, Any]]]]:
    return DATA_CONNECTORS.get(source_type.lower())

# --- File Loader Implementations ---

def load_txt_file(path: str) -> Optional[str]:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return None

def load_pdf_file(path: str) -> Optional[str]:
    try:
        import PyPDF2
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        return None

def load_docx_file(path: str) -> Optional[str]:
    try:
        import docx
        doc = docx.Document(path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception:
        return None

def load_xlsx_file(path: str) -> Optional[str]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True)
        text = []
        # Add sheet names as headers
        for ws in wb.worksheets:
            text.append(f"## Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                text.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
            text.append("\n")  # Add separator between sheets
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error loading Excel file {path}: {e}")
        return None

def load_pptx_file(path: str) -> Optional[str]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text = []
        # Extract slide titles and content
        for i, slide in enumerate(prs.slides):
            text.append(f"## Slide {i+1}: {slide.slide_layout.name}")
            if slide.shapes.title:
                text.append(f"# {slide.shapes.title.text}")
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
            text.append("\n")  # Add separator between slides
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error loading PowerPoint file {path}: {e}")
        return None

def load_csv_file(path: str) -> Optional[str]:
    try:
        import csv
        with open(path, 'r', newline='', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            return "\n".join(["\t".join(row) for row in reader])
    except Exception as e:
        logger.error(f"Error loading CSV file {path}: {e}")
        return None

def load_json_file(path: str) -> Optional[str]:
    try:
        import json
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
            # Pretty print with indentation
            return json.dumps(data, indent=2)
    except Exception as e:
        logger.error(f"Error loading JSON file {path}: {e}")
        return None

def load_html_file(path: str) -> Optional[str]:
    try:
        from bs4 import BeautifulSoup
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            # Extract text from HTML, preserving some structure
            text = []
            # Get title
            if soup.title:
                text.append(f"# {soup.title.string}")
            # Get headings
            for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                level = int(heading.name[1])
                text.append(f"{'#' * level} {heading.get_text().strip()}")
            # Get paragraphs
            for p in soup.find_all('p'):
                text.append(p.get_text().strip())
            # Get list items
            for li in soup.find_all('li'):
                text.append(f"- {li.get_text().strip()}")
            return "\n\n".join([t for t in text if t])
    except Exception as e:
        logger.error(f"Error loading HTML file {path}: {e}")
        # Fallback to plain text loading
        return load_txt_file(path)

# --- Data Connector Stubs ---

def mysql_connector(host: str, user: str, password: str, database: str, query: str) -> List[Dict[str, Any]]:
    """
    Connect to MySQL and return results as file_records.
    Each record: {"path": "mysql://table/row", "name": "table", "content": str(row), ...}
    """
    try:
        import mysql.connector
        conn = mysql.connector.connect(host=host, user=user, password=password, database=database)
        cursor = conn.cursor(dictionary=True)
        cursor.execute(query)
        records = []
        for row in cursor.fetchall():
            records.append({
                "path": f"mysql://{database}/{row.get('id', 'row')}",
                "name": f"{database}",
                "size": len(str(row)),
                "modified": None,
                "content": str(row),
                "source_type": "mysql"
            })
        cursor.close()
        conn.close()
        return records
    except Exception as e:
        return []

def notion_connector(api_key: str, database_id: str) -> List[Dict[str, Any]]:
    # Placeholder: Implement Notion API integration
    return []

def sharepoint_connector(site_url: str, credentials: Dict[str, str]) -> List[Dict[str, Any]]:
    # Placeholder: Implement SharePoint API integration
    return []

def salesforce_connector(credentials: Dict[str, str], soql_query: str) -> List[Dict[str, Any]]:
    # Placeholder: Implement Salesforce API integration
    return []

# Register default loaders
# Text files
for ext in [".txt", ".md", ".toml", ".yaml", ".yml", ".ini", ".cfg", ".conf"]:
    register_file_loader(ext, load_txt_file)

# Code files
for ext in [".py", ".js", ".java", ".c", ".cpp", ".h", ".cs", ".php", ".go", ".rb", ".ts", ".swift"]:
    register_file_loader(ext, load_txt_file)

# Document files
register_file_loader(".pdf", load_pdf_file)
register_file_loader(".docx", load_docx_file)
register_file_loader(".doc", load_docx_file)
register_file_loader(".pptx", load_pptx_file)

# Data files
register_file_loader(".xlsx", load_xlsx_file)
register_file_loader(".csv", load_csv_file)
register_file_loader(".json", load_json_file)
register_file_loader(".html", load_html_file)

# Register data connectors
register_data_connector("mysql", mysql_connector)
register_data_connector("notion", notion_connector)
register_data_connector("sharepoint", sharepoint_connector)
register_data_connector("salesforce", salesforce_connector)

def is_supported_file(filename: str) -> bool:
    ext = os.path.splitext(filename)[1].lower()
    return ext in FILE_LOADERS

def scan_directory(root_dir: str, max_file_size_mb: int = 50) -> List[Dict]:
    """
    Recursively scan root_dir for supported files.

    Args:
        root_dir: Directory to scan
        max_file_size_mb: Maximum file size in MB to process

    Returns:
        List of dicts with file metadata and content
    """
    max_file_size = max_file_size_mb * 1024 * 1024  # Convert MB to bytes
    file_records = []
    root_path = Path(root_dir)

    # Count total files for logging
    total_files = sum(1 for _ in root_path.rglob('*') if _.is_file())
    processed_files = 0
    skipped_files = 0

    logger.info(f"Scanning directory: {root_dir} (found {total_files} total files)")

    for file_path in root_path.rglob('*'):
        if not file_path.is_file():
            continue

        processed_files += 1
        if processed_files % 100 == 0:
            logger.info(f"Processed {processed_files}/{total_files} files")

        fname = file_path.name
        if not is_supported_file(fname):
            skipped_files += 1
            continue

        # Check file size
        try:
            file_size = file_path.stat().st_size
            if file_size > max_file_size:
                logger.warning(f"Skipping {file_path}: file size {file_size/1024/1024:.2f} MB exceeds limit of {max_file_size_mb} MB")
                skipped_files += 1
                continue

            ext = file_path.suffix.lower()
            loader = get_file_loader(ext)
            if loader is None:
                skipped_files += 1
                continue

            content = loader(str(file_path))
            if content is None:
                skipped_files += 1
                continue

            stat = file_path.stat()
            file_records.append({
                "path": str(file_path),
                "name": fname,
                "size": stat.st_size,
                "modified": stat.st_mtime,
                "content": content,
                "source_type": ext.lstrip(".")
            })
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            skipped_files += 1
            continue

    logger.info(f"Scan complete: {len(file_records)} files indexed, {skipped_files} files skipped")
    return file_records

def scan_data_source(source_type: str, **kwargs) -> List[Dict]:
    """
    Scan an external data source (e.g., MySQL, Notion, SharePoint, Salesforce).
    Returns a list of file_records.
    """
    connector = get_data_connector(source_type)
    if connector is None:
        return []
    return connector(**kwargs)
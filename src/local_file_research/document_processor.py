"""
Enhanced document processing module for Local File Deep Research.
"""

import os
import re
import json
import logging
import tempfile
from pathlib import Path
from typing import Dict, List, Any, Optional, Union, BinaryIO, Tuple
from datetime import datetime
import hashlib

# Configure logging
logger = logging.getLogger(__name__)

# --- Document Type Constants ---
TEXT_EXTENSIONS = {".txt", ".md", ".rst", ".log", ".cfg", ".ini", ".conf", ".properties"}
CODE_EXTENSIONS = {
    ".py", ".js", ".java", ".c", ".cpp", ".cs", ".go", ".rb", ".php", ".swift", 
    ".kt", ".ts", ".html", ".css", ".sql", ".sh", ".ps1", ".bat", ".yaml", ".yml",
    ".json", ".xml", ".toml"
}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".doc", ".rtf", ".odt"}
SPREADSHEET_EXTENSIONS = {".xlsx", ".xls", ".csv", ".ods"}
PRESENTATION_EXTENSIONS = {".pptx", ".ppt", ".odp"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg", ".webp"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".flac", ".aac", ".m4a"}
VIDEO_EXTENSIONS = {".mp4", ".avi", ".mov", ".wmv", ".mkv", ".webm"}
ARCHIVE_EXTENSIONS = {".zip", ".tar", ".gz", ".rar", ".7z"}

# --- Document Processing Errors ---
class DocumentProcessingError(Exception):
    """Base exception for document processing errors."""
    pass

class UnsupportedDocumentTypeError(DocumentProcessingError):
    """Exception raised when a document type is not supported."""
    pass

class DocumentTooLargeError(DocumentProcessingError):
    """Exception raised when a document is too large to process."""
    pass

class DocumentParsingError(DocumentProcessingError):
    """Exception raised when a document cannot be parsed."""
    pass

# --- Document Metadata ---
class DocumentMetadata:
    """Class to store and manage document metadata."""
    
    def __init__(self, file_path: Optional[str] = None, content_type: Optional[str] = None):
        """
        Initialize document metadata.
        
        Args:
            file_path: Optional path to the document file
            content_type: Optional content type of the document
        """
        self.metadata = {
            "file_path": file_path,
            "file_name": os.path.basename(file_path) if file_path else None,
            "file_extension": os.path.splitext(file_path)[1].lower() if file_path else None,
            "content_type": content_type,
            "document_type": None,
            "size_bytes": None,
            "created_at": None,
            "modified_at": None,
            "processed_at": datetime.now().isoformat(),
            "word_count": 0,
            "page_count": None,
            "author": None,
            "title": None,
            "language": None,
            "content_hash": None,
            "processing_time_ms": 0,
            "processing_status": "pending",
            "processing_error": None,
            "custom_metadata": {}
        }
        
        # If file path is provided, extract file metadata
        if file_path and os.path.exists(file_path):
            self._extract_file_metadata(file_path)
    
    def _extract_file_metadata(self, file_path: str):
        """
        Extract metadata from a file.
        
        Args:
            file_path: Path to the file
        """
        try:
            stat = os.stat(file_path)
            self.metadata["size_bytes"] = stat.st_size
            self.metadata["created_at"] = datetime.fromtimestamp(stat.st_ctime).isoformat()
            self.metadata["modified_at"] = datetime.fromtimestamp(stat.st_mtime).isoformat()
            
            # Determine document type from extension
            ext = os.path.splitext(file_path)[1].lower()
            if ext in TEXT_EXTENSIONS:
                self.metadata["document_type"] = "text"
            elif ext in CODE_EXTENSIONS:
                self.metadata["document_type"] = "code"
            elif ext in DOCUMENT_EXTENSIONS:
                self.metadata["document_type"] = "document"
            elif ext in SPREADSHEET_EXTENSIONS:
                self.metadata["document_type"] = "spreadsheet"
            elif ext in PRESENTATION_EXTENSIONS:
                self.metadata["document_type"] = "presentation"
            elif ext in IMAGE_EXTENSIONS:
                self.metadata["document_type"] = "image"
            elif ext in AUDIO_EXTENSIONS:
                self.metadata["document_type"] = "audio"
            elif ext in VIDEO_EXTENSIONS:
                self.metadata["document_type"] = "video"
            elif ext in ARCHIVE_EXTENSIONS:
                self.metadata["document_type"] = "archive"
            else:
                self.metadata["document_type"] = "unknown"
        except Exception as e:
            logger.warning(f"Error extracting file metadata: {e}")
    
    def update(self, **kwargs):
        """
        Update metadata fields.
        
        Args:
            **kwargs: Metadata fields to update
        """
        for key, value in kwargs.items():
            if key in self.metadata:
                self.metadata[key] = value
            elif key.startswith("custom_"):
                self.metadata["custom_metadata"][key[7:]] = value
            else:
                self.metadata["custom_metadata"][key] = value
    
    def get(self, key: str, default: Any = None) -> Any:
        """
        Get a metadata field.
        
        Args:
            key: Metadata field name
            default: Default value if field doesn't exist
            
        Returns:
            Metadata field value
        """
        if key in self.metadata:
            return self.metadata[key]
        elif key.startswith("custom_"):
            return self.metadata["custom_metadata"].get(key[7:], default)
        else:
            return self.metadata["custom_metadata"].get(key, default)
    
    def to_dict(self) -> Dict[str, Any]:
        """
        Convert metadata to dictionary.
        
        Returns:
            Metadata dictionary
        """
        return self.metadata
    
    def compute_content_hash(self, content: Union[str, bytes]):
        """
        Compute a hash of the content.
        
        Args:
            content: Document content
        """
        if isinstance(content, str):
            content = content.encode("utf-8")
        
        self.metadata["content_hash"] = hashlib.sha256(content).hexdigest()

# --- Document Processors ---
class BaseDocumentProcessor:
    """Base class for document processors."""
    
    def __init__(self, max_file_size_mb: int = 50):
        """
        Initialize the document processor.
        
        Args:
            max_file_size_mb: Maximum file size in MB
        """
        self.max_file_size_bytes = max_file_size_mb * 1024 * 1024
    
    def process(self, file_path: str) -> Tuple[str, DocumentMetadata]:
        """
        Process a document file.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Tuple of (extracted text, metadata)
            
        Raises:
            DocumentProcessingError: If processing fails
        """
        # Check file size
        if os.path.getsize(file_path) > self.max_file_size_bytes:
            raise DocumentTooLargeError(f"File size exceeds maximum allowed size of {self.max_file_size_bytes // (1024 * 1024)} MB")
        
        # Create metadata
        metadata = DocumentMetadata(file_path)
        
        # Process file based on type
        start_time = datetime.now()
        try:
            content = self._extract_text(file_path, metadata)
            
            # Update metadata
            processing_time = (datetime.now() - start_time).total_seconds() * 1000
            metadata.update(
                processing_time_ms=processing_time,
                processing_status="success",
                word_count=len(content.split())
            )
            
            # Compute content hash
            metadata.compute_content_hash(content)
            
            return content, metadata
        except Exception as e:
            # Update metadata with error
            processing_time = (datetime.now() - start_time).total_seconds() * 1000
            metadata.update(
                processing_time_ms=processing_time,
                processing_status="error",
                processing_error=str(e)
            )
            
            logger.error(f"Error processing document {file_path}: {e}")
            raise DocumentProcessingError(f"Failed to process document: {e}")
    
    def _extract_text(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from a document file.
        
        Args:
            file_path: Path to the document file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            UnsupportedDocumentTypeError: If the document type is not supported
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in TEXT_EXTENSIONS or ext in CODE_EXTENSIONS:
            return self._extract_text_from_text_file(file_path, metadata)
        elif ext in DOCUMENT_EXTENSIONS:
            return self._extract_text_from_document(file_path, metadata)
        elif ext in SPREADSHEET_EXTENSIONS:
            return self._extract_text_from_spreadsheet(file_path, metadata)
        elif ext in PRESENTATION_EXTENSIONS:
            return self._extract_text_from_presentation(file_path, metadata)
        else:
            raise UnsupportedDocumentTypeError(f"Unsupported document type: {ext}")
    
    def _extract_text_from_text_file(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from a text file.
        
        Args:
            file_path: Path to the text file
            metadata: Document metadata
            
        Returns:
            Extracted text
        """
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            
            # Try to detect language
            metadata.update(language=self._detect_language(content))
            
            # Try to extract title from first line
            lines = content.split("\n")
            if lines and lines[0].strip():
                metadata.update(title=lines[0].strip())
            
            return content
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ["latin-1", "cp1252", "iso-8859-1"]:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, read as binary and decode with errors ignored
            with open(file_path, "rb") as f:
                return f.read().decode("utf-8", errors="ignore")
    
    def _extract_text_from_document(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from a document file (PDF, DOCX, etc.).
        
        Args:
            file_path: Path to the document file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the document cannot be parsed
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == ".pdf":
            return self._extract_text_from_pdf(file_path, metadata)
        elif ext in [".docx", ".doc"]:
            return self._extract_text_from_docx(file_path, metadata)
        elif ext == ".rtf":
            return self._extract_text_from_rtf(file_path, metadata)
        elif ext == ".odt":
            return self._extract_text_from_odt(file_path, metadata)
        else:
            raise UnsupportedDocumentTypeError(f"Unsupported document type: {ext}")
    
    def _extract_text_from_pdf(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from a PDF file.
        
        Args:
            file_path: Path to the PDF file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the PDF cannot be parsed
        """
        try:
            import pypdf
        except ImportError:
            try:
                import PyPDF2 as pypdf
            except ImportError:
                raise ImportError("pypdf or PyPDF2 is required to extract text from PDF files")
        
        try:
            text = []
            with open(file_path, "rb") as f:
                pdf = pypdf.PdfReader(f)
                
                # Extract metadata
                info = pdf.metadata
                if info:
                    if hasattr(info, "title") and info.title:
                        metadata.update(title=info.title)
                    if hasattr(info, "author") and info.author:
                        metadata.update(author=info.author)
                
                # Extract text from each page
                metadata.update(page_count=len(pdf.pages))
                for page in pdf.pages:
                    text.append(page.extract_text())
            
            content = "\n\n".join(text)
            
            # Try to detect language
            metadata.update(language=self._detect_language(content))
            
            return content
        except Exception as e:
            raise DocumentParsingError(f"Failed to parse PDF: {e}")
    
    def _extract_text_from_docx(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from a DOCX file.
        
        Args:
            file_path: Path to the DOCX file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the DOCX cannot be parsed
        """
        try:
            import docx
        except ImportError:
            raise ImportError("python-docx is required to extract text from DOCX files")
        
        try:
            doc = docx.Document(file_path)
            
            # Extract metadata
            core_properties = doc.core_properties
            if core_properties:
                if core_properties.title:
                    metadata.update(title=core_properties.title)
                if core_properties.author:
                    metadata.update(author=core_properties.author)
            
            # Extract text from paragraphs
            paragraphs = [p.text for p in doc.paragraphs]
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text)
                    paragraphs.append(" | ".join(row_text))
            
            content = "\n\n".join(paragraphs)
            
            # Try to detect language
            metadata.update(language=self._detect_language(content))
            
            return content
        except Exception as e:
            raise DocumentParsingError(f"Failed to parse DOCX: {e}")
    
    def _extract_text_from_rtf(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from an RTF file.
        
        Args:
            file_path: Path to the RTF file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the RTF cannot be parsed
        """
        try:
            import striprtf
        except ImportError:
            raise ImportError("striprtf is required to extract text from RTF files")
        
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                rtf_text = f.read()
            
            content = striprtf.rtf_to_text(rtf_text)
            
            # Try to detect language
            metadata.update(language=self._detect_language(content))
            
            return content
        except Exception as e:
            raise DocumentParsingError(f"Failed to parse RTF: {e}")
    
    def _extract_text_from_odt(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from an ODT file.
        
        Args:
            file_path: Path to the ODT file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the ODT cannot be parsed
        """
        try:
            import odf.opendocument
            import odf.text
        except ImportError:
            raise ImportError("odfpy is required to extract text from ODT files")
        
        try:
            doc = odf.opendocument.load(file_path)
            
            # Extract metadata
            meta = doc.meta
            if meta:
                title = meta.getElementsByType(odf.text.Title)
                if title and title[0].firstChild:
                    metadata.update(title=title[0].firstChild.data)
                
                creator = meta.getElementsByType(odf.text.Creator)
                if creator and creator[0].firstChild:
                    metadata.update(author=creator[0].firstChild.data)
            
            # Extract text
            paragraphs = []
            for paragraph in doc.getElementsByType(odf.text.P):
                paragraphs.append(paragraph.plainText())
            
            content = "\n\n".join(paragraphs)
            
            # Try to detect language
            metadata.update(language=self._detect_language(content))
            
            return content
        except Exception as e:
            raise DocumentParsingError(f"Failed to parse ODT: {e}")
    
    def _extract_text_from_spreadsheet(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from a spreadsheet file.
        
        Args:
            file_path: Path to the spreadsheet file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the spreadsheet cannot be parsed
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in [".xlsx", ".xls"]:
            return self._extract_text_from_excel(file_path, metadata)
        elif ext == ".csv":
            return self._extract_text_from_csv(file_path, metadata)
        elif ext == ".ods":
            return self._extract_text_from_ods(file_path, metadata)
        else:
            raise UnsupportedDocumentTypeError(f"Unsupported spreadsheet type: {ext}")
    
    def _extract_text_from_excel(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from an Excel file.
        
        Args:
            file_path: Path to the Excel file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the Excel file cannot be parsed
        """
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pandas is required to extract text from Excel files")
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
            
            # Extract text from each sheet
            sheets_text = []
            for sheet_name in sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert to string representation
                sheet_text = f"Sheet: {sheet_name}\n"
                sheet_text += df.to_string(index=False)
                sheets_text.append(sheet_text)
            
            content = "\n\n".join(sheets_text)
            
            # Update metadata
            metadata.update(
                custom_sheet_count=len(sheet_names),
                custom_sheet_names=sheet_names
            )
            
            return content
        except Exception as e:
            raise DocumentParsingError(f"Failed to parse Excel file: {e}")
    
    def _extract_text_from_csv(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from a CSV file.
        
        Args:
            file_path: Path to the CSV file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the CSV file cannot be parsed
        """
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pandas is required to extract text from CSV files")
        
        try:
            # Try different encodings
            encodings = ["utf-8", "latin-1", "cp1252", "iso-8859-1"]
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                # If all encodings fail, use utf-8 with errors ignored
                df = pd.read_csv(file_path, encoding="utf-8", errors="ignore")
            
            # Convert to string representation
            content = df.to_string(index=False)
            
            # Update metadata
            metadata.update(
                custom_row_count=len(df),
                custom_column_count=len(df.columns),
                custom_column_names=list(df.columns)
            )
            
            return content
        except Exception as e:
            raise DocumentParsingError(f"Failed to parse CSV file: {e}")
    
    def _extract_text_from_ods(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from an ODS file.
        
        Args:
            file_path: Path to the ODS file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the ODS file cannot be parsed
        """
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pandas is required to extract text from ODS files")
        
        try:
            # Read all sheets
            df_dict = pd.read_excel(file_path, sheet_name=None, engine="odf")
            sheet_names = list(df_dict.keys())
            
            # Extract text from each sheet
            sheets_text = []
            for sheet_name, df in df_dict.items():
                sheet_text = f"Sheet: {sheet_name}\n"
                sheet_text += df.to_string(index=False)
                sheets_text.append(sheet_text)
            
            content = "\n\n".join(sheets_text)
            
            # Update metadata
            metadata.update(
                custom_sheet_count=len(sheet_names),
                custom_sheet_names=sheet_names
            )
            
            return content
        except Exception as e:
            raise DocumentParsingError(f"Failed to parse ODS file: {e}")
    
    def _extract_text_from_presentation(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from a presentation file.
        
        Args:
            file_path: Path to the presentation file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the presentation cannot be parsed
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in [".pptx", ".ppt"]:
            return self._extract_text_from_powerpoint(file_path, metadata)
        elif ext == ".odp":
            return self._extract_text_from_odp(file_path, metadata)
        else:
            raise UnsupportedDocumentTypeError(f"Unsupported presentation type: {ext}")
    
    def _extract_text_from_powerpoint(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from a PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the PowerPoint file cannot be parsed
        """
        try:
            import pptx
        except ImportError:
            raise ImportError("python-pptx is required to extract text from PowerPoint files")
        
        try:
            presentation = pptx.Presentation(file_path)
            
            # Extract metadata
            core_properties = presentation.core_properties
            if core_properties:
                if core_properties.title:
                    metadata.update(title=core_properties.title)
                if core_properties.author:
                    metadata.update(author=core_properties.author)
            
            # Extract text from slides
            slides_text = []
            for i, slide in enumerate(presentation.slides, 1):
                slide_text = f"Slide {i}:\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                slides_text.append(slide_text)
            
            content = "\n\n".join(slides_text)
            
            # Update metadata
            metadata.update(
                page_count=len(presentation.slides),
                custom_slide_count=len(presentation.slides)
            )
            
            # Try to detect language
            metadata.update(language=self._detect_language(content))
            
            return content
        except Exception as e:
            raise DocumentParsingError(f"Failed to parse PowerPoint file: {e}")
    
    def _extract_text_from_odp(self, file_path: str, metadata: DocumentMetadata) -> str:
        """
        Extract text from an ODP file.
        
        Args:
            file_path: Path to the ODP file
            metadata: Document metadata
            
        Returns:
            Extracted text
            
        Raises:
            DocumentParsingError: If the ODP file cannot be parsed
        """
        try:
            import odf.opendocument
            import odf.text
        except ImportError:
            raise ImportError("odfpy is required to extract text from ODP files")
        
        try:
            doc = odf.opendocument.load(file_path)
            
            # Extract metadata
            meta = doc.meta
            if meta:
                title = meta.getElementsByType(odf.text.Title)
                if title and title[0].firstChild:
                    metadata.update(title=title[0].firstChild.data)
                
                creator = meta.getElementsByType(odf.text.Creator)
                if creator and creator[0].firstChild:
                    metadata.update(author=creator[0].firstChild.data)
            
            # Extract text from paragraphs
            paragraphs = []
            for paragraph in doc.getElementsByType(odf.text.P):
                paragraphs.append(paragraph.plainText())
            
            content = "\n\n".join(paragraphs)
            
            # Try to detect language
            metadata.update(language=self._detect_language(content))
            
            return content
        except Exception as e:
            raise DocumentParsingError(f"Failed to parse ODP file: {e}")
    
    def _detect_language(self, text: str) -> Optional[str]:
        """
        Detect the language of a text.
        
        Args:
            text: Text to detect language for
            
        Returns:
            Detected language code or None
        """
        try:
            import langdetect
            
            # Use a sample of the text for faster detection
            sample = text[:5000]
            if not sample:
                return None
            
            return langdetect.detect(sample)
        except ImportError:
            logger.warning("langdetect is not installed. Language detection is disabled.")
            return None
        except Exception:
            return None

# --- Document Processor Factory ---
class DocumentProcessorFactory:
    """Factory for creating document processors."""
    
    @staticmethod
    def create_processor(max_file_size_mb: int = 50) -> BaseDocumentProcessor:
        """
        Create a document processor.
        
        Args:
            max_file_size_mb: Maximum file size in MB
            
        Returns:
            Document processor
        """
        return BaseDocumentProcessor(max_file_size_mb=max_file_size_mb)

# --- Main Processing Function ---
def process_document(file_path: str, max_file_size_mb: int = 50) -> Tuple[str, Dict[str, Any]]:
    """
    Process a document file.
    
    Args:
        file_path: Path to the document file
        max_file_size_mb: Maximum file size in MB
        
    Returns:
        Tuple of (extracted text, metadata dictionary)
        
    Raises:
        DocumentProcessingError: If processing fails
    """
    processor = DocumentProcessorFactory.create_processor(max_file_size_mb=max_file_size_mb)
    content, metadata = processor.process(file_path)
    return content, metadata.to_dict()

def process_document_from_bytes(content: bytes, filename: str, max_file_size_mb: int = 50) -> Tuple[str, Dict[str, Any]]:
    """
    Process a document from bytes.
    
    Args:
        content: Document content as bytes
        filename: Original filename
        max_file_size_mb: Maximum file size in MB
        
    Returns:
        Tuple of (extracted text, metadata dictionary)
        
    Raises:
        DocumentProcessingError: If processing fails
    """
    # Check file size
    if len(content) > max_file_size_mb * 1024 * 1024:
        raise DocumentTooLargeError(f"File size exceeds maximum allowed size of {max_file_size_mb} MB")
    
    # Save to temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as temp_file:
        temp_file.write(content)
        temp_path = temp_file.name
    
    try:
        # Process the temporary file
        return process_document(temp_path, max_file_size_mb=max_file_size_mb)
    finally:
        # Clean up temporary file
        try:
            os.unlink(temp_path)
        except:
            pass

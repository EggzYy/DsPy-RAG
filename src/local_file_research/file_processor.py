"""
File processor module for handling different file types.
"""

# Import the logging silencer first to suppress all logging errors
try:
    from .logging_silence import silence_all_logging_errors
except ImportError:
    try:
        from src.local_file_research.logging_silence import silence_all_logging_errors
    except ImportError:
        pass  # Silently continue if the module is not available

import logging
import os
import io
import base64
from typing import Optional, Dict, Any, Tuple

# Configure logging
logger = logging.getLogger(__name__)

# Try to import specialized file processing libraries
try:
    import fitz  # PyMuPDF for PDF processing
    PYMUPDF_AVAILABLE = True
    logger.info("PyMuPDF is available for PDF processing")
except ImportError:
    PYMUPDF_AVAILABLE = False
    logger.warning("PyMuPDF is not available, PDF text extraction will be limited")

try:
    import docx
    DOCX_AVAILABLE = True
    logger.info("python-docx is available for DOCX processing")
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx is not available, DOCX text extraction will be limited")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
    logger.info("python-pptx is available for PPTX processing")
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx is not available, PPTX text extraction will be limited")

def process_file(file_bytes: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
    """
    Process a file based on its type and extract text and metadata.

    Args:
        file_bytes: The binary content of the file
        filename: The name of the file

    Returns:
        Tuple containing:
            - Extracted text content
            - Metadata dictionary
    """
    # Get file extension
    _, ext = os.path.splitext(filename)
    ext = ext.lower().lstrip('.')

    metadata = {
        "filename": filename,
        "size_bytes": len(file_bytes),
        "file_type": ext,
        "extraction_method": "unknown"
    }

    # Process based on file type
    if ext == 'pdf':
        return process_pdf(file_bytes, metadata)
    elif ext in ['docx', 'doc']:
        return process_docx(file_bytes, metadata)
    elif ext in ['pptx', 'ppt']:
        return process_pptx(file_bytes, metadata)
    elif ext in ['txt', 'md', 'py', 'js', 'html', 'css', 'json', 'xml', 'csv']:
        return process_text(file_bytes, metadata)
    else:
        # Default processing for unknown file types
        return process_binary(file_bytes, metadata)

def process_pdf(file_bytes: bytes, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
    """Process PDF files using PyMuPDF if available."""
    if PYMUPDF_AVAILABLE:
        try:
            # Use PyMuPDF to extract text
            with fitz.open(stream=file_bytes, filetype="pdf") as pdf_document:
                text = ""
                page_count = len(pdf_document)
                metadata["page_count"] = page_count

                # Extract text from each page
                for page_num in range(page_count):
                    page = pdf_document[page_num]
                    text += page.get_text() + "\n\n"

                # Get document info
                info = pdf_document.metadata
                if info:
                    metadata["title"] = info.get("title", "")
                    metadata["author"] = info.get("author", "")
                    metadata["subject"] = info.get("subject", "")
                    metadata["keywords"] = info.get("keywords", "")

                metadata["extraction_method"] = "pymupdf"
                logger.info(f"Successfully extracted {len(text)} characters from PDF using PyMuPDF")
                return text, metadata
        except Exception as e:
            logger.error(f"Error extracting text from PDF with PyMuPDF: {e}")

    # Fallback: Return base64 encoded content
    logger.warning("Using fallback method for PDF processing")
    metadata["extraction_method"] = "base64"
    return base64.b64encode(file_bytes).decode("utf-8"), metadata

def process_docx(file_bytes: bytes, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
    """Process DOCX files using python-docx if available."""
    if DOCX_AVAILABLE:
        try:
            # Use python-docx to extract text
            doc = docx.Document(io.BytesIO(file_bytes))

            # Extract text with proper Unicode handling
            paragraphs = []
            for paragraph in doc.paragraphs:
                try:
                    # Handle each paragraph individually to isolate potential encoding issues
                    paragraphs.append(paragraph.text)
                except Exception as para_error:
                    logger.warning(f"Error extracting paragraph text: {para_error}")
                    # Add placeholder for problematic paragraph
                    paragraphs.append("[Error: Could not extract paragraph text]")

            # Join paragraphs with proper Unicode handling
            text = "\n\n".join(paragraphs)

            # Get document properties
            metadata["paragraph_count"] = len(doc.paragraphs)
            metadata["extraction_method"] = "python-docx"

            logger.info(f"Successfully extracted {len(text)} characters from DOCX using python-docx")
            return text, metadata
        except Exception as e:
            logger.error(f"Error extracting text from DOCX with python-docx: {e}")

    # Fallback: Try multiple encoding approaches
    encoding_attempts = ["utf-8", "latin-1", "cp1252", "utf-16"]
    for encoding in encoding_attempts:
        try:
            text = file_bytes.decode(encoding)
            metadata["extraction_method"] = f"decoded-{encoding}"
            logger.info(f"Successfully decoded DOCX content using {encoding} encoding")
            return text, metadata
        except UnicodeDecodeError:
            continue

    # If all decoding attempts fail, return base64 encoded content
    logger.warning("All decoding attempts failed. Using base64 fallback for DOCX processing")
    metadata["extraction_method"] = "base64"
    return base64.b64encode(file_bytes).decode("utf-8"), metadata

def process_pptx(file_bytes: bytes, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
    """Process PPTX files using python-pptx if available."""
    if PPTX_AVAILABLE:
        try:
            # Use python-pptx to extract text
            prs = Presentation(io.BytesIO(file_bytes))
            text = ""

            # Extract text from slides
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n\n"

            # Get presentation properties
            metadata["slide_count"] = len(prs.slides)
            metadata["extraction_method"] = "python-pptx"

            logger.info(f"Successfully extracted {len(text)} characters from PPTX using python-pptx")
            return text, metadata
        except Exception as e:
            logger.error(f"Error extracting text from PPTX with python-pptx: {e}")

    # Fallback: Return base64 encoded content
    logger.warning("Using fallback method for PPTX processing")
    metadata["extraction_method"] = "base64"
    return base64.b64encode(file_bytes).decode("utf-8"), metadata

def process_text(file_bytes: bytes, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
    """Process text files."""
    try:
        # Try UTF-8 decoding first
        text = file_bytes.decode("utf-8")
        metadata["extraction_method"] = "utf8"
        return text, metadata
    except UnicodeDecodeError:
        try:
            # Try Latin-1 as fallback
            text = file_bytes.decode("latin-1")
            metadata["extraction_method"] = "latin1"
            return text, metadata
        except Exception as e:
            logger.error(f"Error decoding text file: {e}")
            # If all decoding fails, return base64 encoded content
            metadata["extraction_method"] = "base64"
            return base64.b64encode(file_bytes).decode("utf-8"), metadata

def process_binary(file_bytes: bytes, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
    """Process binary files by returning base64 encoded content."""
    metadata["extraction_method"] = "base64"
    return base64.b64encode(file_bytes).decode("utf-8"), metadata
